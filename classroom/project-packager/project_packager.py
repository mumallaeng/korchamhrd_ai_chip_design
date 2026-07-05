#!/usr/bin/env python3
from __future__ import annotations

import argparse
import copy
import glob
import json
import os
import re
import shutil
import sys
from pathlib import Path
from typing import Any, NoReturn


def fail(message: str) -> NoReturn:
    raise SystemExit(f"error: {message}")


def load_json(path: Path) -> Any:
    with path.open("r", encoding="utf-8") as handle:
        return json.load(handle)


def load_text(path: Path) -> str:
    return path.read_text(encoding="utf-8")


def resolve_path(base_dir: Path, raw_path: str) -> Path:
    path = Path(os.path.expanduser(raw_path))
    if path.is_absolute():
        return path
    return (base_dir / path).resolve()


def ensure_parent(path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)


def remove_existing(path: Path) -> None:
    if path.is_symlink() or path.is_file():
        path.unlink()
    elif path.is_dir():
        shutil.rmtree(path)


def render_string(template: str, variables: dict[str, str]) -> str:
    try:
        return template.format_map(variables)
    except KeyError as exc:
        fail(f"missing variable '{exc.args[0]}' while rendering: {template}")


def render_data(value: Any, variables: dict[str, str]) -> Any:
    if isinstance(value, str):
        return render_string(value, variables)
    if isinstance(value, list):
        return [render_data(item, variables) for item in value]
    if isinstance(value, dict):
        return {key: render_data(item, variables) for key, item in value.items()}
    return value


def normalize_whitespace(text: str) -> str:
    return " ".join(text.split())


def resolve_variables(config: dict[str, Any], config_dir: Path, overrides: dict[str, str]) -> dict[str, str]:
    raw_variables = copy.deepcopy(config.get("variables", {}))
    raw_variables.update(overrides)
    variables: dict[str, str] = {"config_dir": str(config_dir)}

    for _ in range(20):
        changed = False
        for key, raw_value in raw_variables.items():
            candidate = str(raw_value)
            rendered = render_string(candidate, {**variables, **raw_variables})
            if variables.get(key) != rendered:
                variables[key] = rendered
                changed = True
        if not changed:
            break
    else:
        fail("variables did not converge after 20 render passes")

    return variables


def prepare_target_dir(target_dir: Path) -> None:
    target_dir.mkdir(parents=True, exist_ok=True)


def replace_file(source: Path, target: Path) -> None:
    if not source.is_file():
        fail(f"missing source file: {source}")
    ensure_parent(target)
    remove_existing(target)
    shutil.copy2(source, target)


def replace_directory(source: Path, target: Path) -> None:
    if not source.is_dir():
        fail(f"missing source directory: {source}")
    ensure_parent(target)
    remove_existing(target)
    shutil.copytree(source, target)


def write_text_file(target: Path, content: str) -> None:
    ensure_parent(target)
    target.write_text(content, encoding="utf-8")


def load_config_block(
    section: dict[str, Any],
    key_inline: str,
    key_file: str,
    config_dir: Path,
    variables: dict[str, str],
    default: Any,
) -> Any:
    if key_file in section:
        file_path = resolve_path(config_dir, render_string(section[key_file], variables))
        data = load_json(file_path)
        return render_data(data, variables)
    if key_inline in section:
        return render_data(section[key_inline], variables)
    return default


def apply_regex_replacements(config: dict[str, Any], config_dir: Path, variables: dict[str, str]) -> None:
    for item in config.get("regex_replacements", []):
        rendered = render_data(item, variables)
        path = resolve_path(config_dir, rendered["path"])
        pattern = rendered["pattern"]
        replacement = rendered.get("replacement", "")
        flags = 0
        for flag_name in rendered.get("flags", []):
            if flag_name == "MULTILINE":
                flags |= re.MULTILINE
            elif flag_name == "DOTALL":
                flags |= re.DOTALL
            else:
                fail(f"unsupported regex flag: {flag_name}")
        count = int(rendered.get("count", 0))
        original = path.read_text(encoding="utf-8")
        updated, replacements = re.subn(pattern, replacement, original, count=count, flags=flags)
        if rendered.get("require_match", True) and replacements == 0:
            fail(f"regex did not match in {path}: {pattern}")
        path.write_text(updated, encoding="utf-8")


def apply_delete_globs(config: dict[str, Any], config_dir: Path, variables: dict[str, str]) -> None:
    for raw_pattern in config.get("delete_globs", []):
        rendered_pattern = render_string(raw_pattern, variables)
        pattern_path = Path(os.path.expanduser(rendered_pattern))
        if pattern_path.is_absolute():
            pattern = str(pattern_path)
        else:
            pattern = str((config_dir / pattern_path).resolve())
        matches = glob.glob(pattern, recursive=True)
        for match in matches:
            path = resolve_path(config_dir, match)
            if path.exists() or path.is_symlink():
                remove_existing(path)


def build_text_files(config: dict[str, Any], config_dir: Path, variables: dict[str, str]) -> None:
    for item in config.get("text_files", []):
        rendered = render_data(item, variables)
        target = resolve_path(config_dir, rendered["target"])
        if "content" in rendered:
            content = rendered["content"]
        elif "source_markdown" in rendered:
            content = render_string(load_text(resolve_path(config_dir, rendered["source_markdown"])), variables)
        elif "source_text" in rendered:
            content = render_string(load_text(resolve_path(config_dir, rendered["source_text"])), variables)
        else:
            fail(f"text file entry needs content or source file: {rendered}")
        write_text_file(target, content)


def build_copies(config: dict[str, Any], config_dir: Path, variables: dict[str, str]) -> None:
    for item in config.get("copies", []):
        rendered = render_data(item, variables)
        source = resolve_path(config_dir, rendered["source"])
        target = resolve_path(config_dir, rendered["target"])
        replace_file(source, target)


def build_directories(config: dict[str, Any], config_dir: Path, variables: dict[str, str]) -> None:
    for item in config.get("directories", []):
        rendered = render_data(item, variables)
        source = resolve_path(config_dir, rendered["source"])
        target = resolve_path(config_dir, rendered["target"])
        replace_directory(source, target)


def apply_journal_section(config: dict[str, Any], config_dir: Path, variables: dict[str, str]) -> None:
    section = config.get("journal_md")
    if not section:
        return

    rendered = render_data(section, variables)
    target = resolve_path(config_dir, rendered["target"])
    if "content" in rendered:
        content = rendered["content"]
    elif "source_markdown" in rendered:
        content = render_string(load_text(resolve_path(config_dir, rendered["source_markdown"])), variables)
    elif "source_text" in rendered:
        content = render_string(load_text(resolve_path(config_dir, rendered["source_text"])), variables)
    else:
        fail("journal_md requires content or source_markdown/source_text")
    write_text_file(target, content)


def import_docx():
    try:
        from docx import Document  # type: ignore
    except ImportError:
        fail("python-docx is not installed. Run ./setup.sh first.")
    return Document


def import_openpyxl():
    try:
        from openpyxl import load_workbook  # type: ignore
    except ImportError:
        fail("openpyxl is not installed. Run ./setup.sh first.")
    return load_workbook


def import_pptx():
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        fail("python-pptx is not installed. Run ./setup.sh first.")
    return Presentation


def apply_docx_section(config: dict[str, Any], config_dir: Path, variables: dict[str, str]) -> None:
    section = config.get("report_docx")
    if not section:
        return

    rendered = render_data(section, variables)
    source_key = "template" if "template" in rendered else "source"
    if source_key not in rendered:
        fail("report_docx requires either template or source")

    source = resolve_path(config_dir, rendered[source_key])
    target = resolve_path(config_dir, rendered["target"])
    replace_file(source, target)

    paragraph_updates = load_config_block(rendered, "paragraph_updates", "paragraph_updates_file", config_dir, variables, [])
    table_updates = load_config_block(rendered, "table_updates", "tables_file", config_dir, variables, [])
    append_paragraphs = rendered.get("append_paragraphs", [])

    if not paragraph_updates and not table_updates and not append_paragraphs:
        return

    Document = import_docx()
    document = Document(str(target))

    for update in paragraph_updates:
        index = int(update["index"])
        if index < 0 or index >= len(document.paragraphs):
            fail(f"paragraph index out of range: {index}")
        document.paragraphs[index].text = update["text"]

    for update in table_updates:
        table_index = int(update["table"])
        row_index = int(update["row"])
        col_index = int(update["col"])
        try:
            cell = document.tables[table_index].rows[row_index].cells[col_index]
        except IndexError:
            fail(f"table cell out of range: {update}")
        cell.text = update["text"]

    for text in append_paragraphs:
        document.add_paragraph(text)

    document.save(str(target))


def merged_anchor(ws, cell_ref: str) -> str:
    for merged_range in ws.merged_cells.ranges:
        if cell_ref in merged_range:
            return merged_range.start_cell.coordinate
    return cell_ref


def set_cell_value(ws, cell_ref: str, value: Any) -> None:
    anchor_ref = merged_anchor(ws, cell_ref)
    ws[anchor_ref] = value


def apply_merge_updates(ws, updates: list[dict[str, Any]]) -> None:
    for update in updates:
        action = update["action"]
        cell_range = update["range"]
        if action == "remove":
            ws.unmerge_cells(cell_range)
        elif action == "add":
            ws.merge_cells(cell_range)
        else:
            fail(f"unsupported merge action: {action}")


def prune_workbook_sheets(workbook, rendered: dict[str, Any]) -> None:
    keep_sheets = rendered.get("keep_sheets")
    remove_sheets = rendered.get("remove_sheets")

    if keep_sheets and remove_sheets:
        fail("schedule_xlsx cannot use keep_sheets and remove_sheets together")

    if keep_sheets:
        keep_set = set(keep_sheets)
        for name in keep_sheets:
            if name not in workbook.sheetnames:
                fail(f"schedule_xlsx keep_sheets references missing sheet: {name}")
        for sheet_name in list(workbook.sheetnames):
            if sheet_name not in keep_set:
                workbook.remove(workbook[sheet_name])
        workbook.active = workbook.sheetnames.index(keep_sheets[0])
        return

    if remove_sheets:
        for sheet_name in remove_sheets:
            if sheet_name not in workbook.sheetnames:
                continue
            if len(workbook.sheetnames) == 1:
                fail("schedule_xlsx cannot remove the last remaining sheet")
            workbook.remove(workbook[sheet_name])


def apply_xlsx_section(config: dict[str, Any], config_dir: Path, variables: dict[str, str]) -> None:
    section = config.get("schedule_xlsx")
    if not section:
        return

    rendered = render_data(section, variables)
    source_key = "template" if "template" in rendered else "source"
    if source_key not in rendered:
        fail("schedule_xlsx requires either template or source")

    source = resolve_path(config_dir, rendered[source_key])
    target = resolve_path(config_dir, rendered["target"])
    replace_file(source, target)

    cell_updates = load_config_block(rendered, "cell_updates", "cell_updates_file", config_dir, variables, [])
    row_updates = load_config_block(rendered, "row_updates", "rows_file", config_dir, variables, [])
    merge_updates = load_config_block(rendered, "merge_updates", "merge_updates_file", config_dir, variables, [])

    if not cell_updates and not row_updates and not merge_updates:
        return

    load_workbook = import_openpyxl()
    workbook = load_workbook(str(target))
    prune_workbook_sheets(workbook, rendered)

    for update in merge_updates:
        sheet_name = update.get("sheet") or workbook.active.title
        ws = workbook[sheet_name]
        apply_merge_updates(ws, [update])

    for update in cell_updates:
        sheet_name = update.get("sheet") or rendered.get("sheet") or workbook.active.title
        ws = workbook[sheet_name]
        set_cell_value(ws, update["cell"], update.get("value", ""))

    for update in row_updates:
        sheet_name = update.get("sheet") or rendered.get("sheet") or workbook.active.title
        ws = workbook[sheet_name]
        row_index = int(update["row"])
        values = update.get("values", {})
        for column, value in values.items():
            set_cell_value(ws, f"{column}{row_index}", value)

    workbook.save(str(target))


def inspect_docx(path: Path) -> str:
    Document = import_docx()
    document = Document(str(path))
    lines: list[str] = []
    lines.append(f"Path: {path}")
    lines.append(f"Paragraphs: {len(document.paragraphs)}")
    lines.append("")
    lines.append("Paragraph indices")
    for index, paragraph in enumerate(document.paragraphs):
        text = normalize_whitespace(paragraph.text)
        lines.append(f"[{index:03d}] {text}")
    lines.append("")
    lines.append(f"Tables: {len(document.tables)}")
    for table_index, table in enumerate(document.tables):
        lines.append("")
        lines.append(f"Table {table_index}")
        for row_index, row in enumerate(table.rows):
            cells = [normalize_whitespace(cell.text) for cell in row.cells]
            lines.append(f"[{row_index:03d}] " + " | ".join(cells))
    return "\n".join(lines)


def inspect_xlsx(path: Path) -> str:
    load_workbook = import_openpyxl()
    workbook = load_workbook(str(path))
    lines: list[str] = []
    lines.append(f"Path: {path}")
    lines.append(f"Sheets: {', '.join(workbook.sheetnames)}")
    for ws in workbook.worksheets:
        lines.append("")
        lines.append(f"Sheet: {ws.title}")
        lines.append(f"Dimensions: rows={ws.max_row}, cols={ws.max_column}")
        merged_ranges = [str(item) for item in ws.merged_cells.ranges]
        lines.append(f"Merged ranges ({len(merged_ranges)}):")
        for merged_range in merged_ranges:
            lines.append(f"  {merged_range}")
        lines.append("Non-empty rows:")
        for row_index in range(1, ws.max_row + 1):
            values = []
            for cell in ws[row_index]:
                value = cell.value
                if value is not None and str(value).strip() != "":
                    values.append(f"{cell.coordinate}={value}")
            if values:
                lines.append(f"  row {row_index}: " + ", ".join(values))
    return "\n".join(lines)


def outline_pptx(path: Path) -> str:
    Presentation = import_pptx()
    presentation = Presentation(str(path))
    lines: list[str] = []
    lines.append(f"# {path.name}")
    lines.append("")
    for index, slide in enumerate(presentation.slides, start=1):
        lines.append(f"## Slide {index}")
        slide_lines: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = normalize_whitespace(shape.text)
            if text:
                slide_lines.append(text)
        if slide_lines:
            for text in slide_lines:
                lines.append(f"- {text}")
        else:
            lines.append("- (no text)")
        lines.append("")
    return "\n".join(lines).rstrip() + "\n"


def write_or_print(text: str, output_path: str | None) -> None:
    if output_path:
        target = Path(os.path.expanduser(output_path))
        ensure_parent(target)
        target.write_text(text, encoding="utf-8")
        print(target)
        return
    print(text)


def build_package(config_path: Path, args: argparse.Namespace) -> None:
    config = load_json(config_path)
    config_dir = config_path.parent.resolve()

    overrides = {}
    for raw_pair in args.var or []:
        if "=" not in raw_pair:
            fail(f"invalid --var format, expected KEY=VALUE: {raw_pair}")
        key, value = raw_pair.split("=", 1)
        overrides[key] = value

    variables = resolve_variables(config, config_dir, overrides)

    raw_target_dir = args.target_dir or config.get("target_dir")
    if not raw_target_dir:
        fail("config must define target_dir or you must pass --target-dir")
    target_dir = resolve_path(config_dir, render_string(raw_target_dir, variables))
    variables["target_dir"] = str(target_dir)

    prepare_target_dir(target_dir)
    build_copies(config, config_dir, variables)
    build_directories(config, config_dir, variables)
    build_text_files(config, config_dir, variables)
    apply_journal_section(config, config_dir, variables)
    apply_docx_section(config, config_dir, variables)
    apply_xlsx_section(config, config_dir, variables)
    apply_regex_replacements(config, config_dir, variables)
    apply_delete_globs(config, config_dir, variables)
    print(target_dir)


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Package repeatable classroom project submission folders.")
    subparsers = parser.add_subparsers(dest="command", required=True)

    build_cmd = subparsers.add_parser("build", help="Build one package from a JSON config.")
    build_cmd.add_argument("config", help="Path to the package JSON config.")
    build_cmd.add_argument("--target-dir", help="Override output directory.")
    build_cmd.add_argument("--var", action="append", default=[], help="Override variables as KEY=VALUE.")

    inspect_docx_cmd = subparsers.add_parser("inspect-docx", help="Print paragraph and table indices for a DOCX file.")
    inspect_docx_cmd.add_argument("path", help="DOCX template or output path.")
    inspect_docx_cmd.add_argument("--output", help="Optional text output file.")

    inspect_xlsx_cmd = subparsers.add_parser("inspect-xlsx", help="Print sheet, merge, and non-empty cell layout for an XLSX file.")
    inspect_xlsx_cmd.add_argument("path", help="XLSX template or output path.")
    inspect_xlsx_cmd.add_argument("--output", help="Optional text output file.")

    outline_pptx_cmd = subparsers.add_parser("outline-pptx", help="Extract per-slide text outline from a PPTX file.")
    outline_pptx_cmd.add_argument("path", help="PPTX presentation path.")
    outline_pptx_cmd.add_argument("--output", help="Optional markdown output file.")

    return parser


def main(argv: list[str] | None = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)

    if args.command == "build":
        build_package(resolve_path(Path.cwd(), args.config), args)
        return 0
    if args.command == "inspect-docx":
        text = inspect_docx(resolve_path(Path.cwd(), args.path))
        write_or_print(text, args.output)
        return 0
    if args.command == "inspect-xlsx":
        text = inspect_xlsx(resolve_path(Path.cwd(), args.path))
        write_or_print(text, args.output)
        return 0
    if args.command == "outline-pptx":
        text = outline_pptx(resolve_path(Path.cwd(), args.path))
        write_or_print(text, args.output)
        return 0

    fail(f"unsupported command: {args.command}")


if __name__ == "__main__":
    sys.exit(main())
